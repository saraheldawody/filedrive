from django.shortcuts import render
from django.shortcuts import redirect
from .models import User, File, LookupWords
from passlib.context import CryptContext
import PyPDF2 
import io
from pptx import Presentation
import nltk
from nltk.corpus import stopwords
from nltk.stem import PorterStemmer 
import string
import validators
from django.db.models import Q
pwd_context = CryptContext(
        schemes=["pbkdf2_sha256"],
        default="pbkdf2_sha256",
        pbkdf2_sha256__default_rounds=30000
)
def encrypt_password(password):
    return pwd_context.encrypt(password)


def check_encrypted_password(password, hashed):
    return pwd_context.verify(password, hashed)
stop_words = set(stopwords.words('english')) 
def clean_line(line):
    line = line.replace('\t',' ')
    line = line.replace('\n',' ')
    line = line.replace('\r',' ')
    line = line.replace('\a',' ')
    line = line.replace('\f',' ')
    line = line.replace('\v',' ')
    line = line.replace('\b',' ')
    return line.split(' ')
def clean_word(strng):
    exclude = set(string.punctuation)
    valid=validators.url(strng)
    if valid:
        return None
    if len(strng)< 2:
        return None
    if (strng == ' ')or (strng == ''):
        return None
    strng = ''.join(ch for ch in strng if ch not in exclude)
    strng = ''.join([i for i in strng if not i.isdigit()])
    ps = PorterStemmer() 
    strng = ps.stem(strng)
    if strng.lower() in stop_words:
        return None
    return strng 
                        

def index(request):
    context={}
    if 'loggedIn' in request.session:
        if request.session['loggedIn']:
            context['loggedin'] = True
    return render(request, 'index.html',context)

def signin(request):
    if request.method == "POST":
        if User.objects.filter(Email=request.POST.get("Email" , "")).count()>0:
            user = User.objects.get(Email=request.POST.get("Email" , ""))
            print(check_encrypted_password(request.POST.get("Password" , ""),user.Password))
            print(request.POST.get("Password" , ""))
            if check_encrypted_password(request.POST.get("Password" , ""),user.Password):
                request.session['loggedIn'] = True
                request.session['username'] = user.Username
                request.session['email']=user.Email
                return redirect ('../')
            else:
                context = {
                'error':"The Password is incorrect",
                'email':request.POST.get("Email" , "")
                }
                return render(request, 'signin.html',context)
        else:
            context = {
                'error':"This Email is not exist",
                'email':request.POST.get("Email" , "")
            }
            return render(request, 'signin.html',context)
    else:
        context={}
        return render(request, 'signin.html',context)

def register(request):
    if request.method == "POST":
        newUser = User.objects.create(
            FirstName=request.POST.get("FirstName" , ""),
            LastName=request.POST.get("LastName" , ""),
            Email=request.POST.get("Email" , ""),
            Username=request.POST.get("Username" , ""),
            Password=encrypt_password(request.POST.get("Password" , "")))
        newUser.save()
        return redirect('../done/')
    else:
        context={}
        return render(request, 'register.html',context)

def signout(request):
    request.session['loggedIn'] = False
    request.session['username'] = ''
    request.session['email']= ''
    return redirect('/')

def delete(request,id):
    f = File.objects.get(id=id)
    f.delete()
    return redirect('../files/')

def files(request):
    if request.method == "POST":
        text = request.POST.get("search","")
        words = text.split(' ')
        values = [w for w in words if w not in stop_words]
        words_cleaned = [clean_word(w) for w in values]
        user = User.objects.get(Email=request.session['email'])
        fs = File.objects.filter(User=user)
        result_content = LookupWords.objects.filter(File__in=fs).filter(Word__in = words_cleaned)
        query = Q()  # empty Q object
        for word in values:
            # 'or' the queries together
            query |= Q(Title__icontains=word)
        result_title = fs.filter(query)
        query = Q()  # empty Q object
        for word in values:
            # 'or' the queries together
            query |= Q(Description__icontains=word)
        result_desc = fs.filter(query)

        files=[]
        for s in result_content:
            if s.File in files:
                continue
            else:
                files.append(s.File)
        for s in result_title:
            if s in files:
                continue
            else:
                files.append(s)
        for s in result_desc:
            if s in files:
                continue
            else:
                files.append(s)
        context={
            'files': files,
            'search':text
        }
        return render(request, 'files.html',context)
    else:
        user = User.objects.get(Email=request.session['email'])
        files = File.objects.filter(User=user)
        context={
            'files': files
        }
        return render(request, 'files.html',context)

def regdone(request):
    context={}
    return render(request, 'regdone.html',context)

def upload(request):
    if request.method == "POST":
        uploadedFile = request.FILES['file']
        name = uploadedFile.name
        FileObj = uploadedFile.read()
        print("filename is ",name)
        content = []
        if name.endswith('.pdf'):
            try:
                print ("extension is pdf")
                pdfReader = PyPDF2.PdfFileReader(io.BytesIO(FileObj))
                NumPages = pdfReader.numPages
                i = 0
                while (i<NumPages):
                    text = pdfReader.getPage(i)
                    content.append(text.extractText())
                    i +=1
                print(content)
            except Exception:
                context={
                    'err':'This PDF is corrupted or cannot be uploaded please try again or try another file.'
                    }
                return render(request, 'upload.html',context)
        elif name.endswith('.pptx'):
            prs = Presentation(io.BytesIO(FileObj))
            print("----------------------")
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content.append(shape.text)
            print(content)
        else:
            print("invalid file type")
        strings_dict={}
        
        if content != []:
            for text in content:
                lines = text.split('\n')
                for line in lines:
                    if (line == '') or (line == ' '):
                        continue
                    strings = clean_line(line)
                    for strng in strings:
                        strng = clean_word(strng)
                        if strng == None:
                            continue
                        if strng.lower() in strings_dict.keys():
                            strings_dict[strng.lower()]+=1
                        else:
                            strings_dict[strng.lower()]=1
        title = request.POST.get('title','')
        description = request.POST.get('description','')
        user = User.objects.get(Email=request.session['email'])
        newFile = File.objects.create(
            Title=title,
            Description=description,
            File=uploadedFile,
            User=user
        )
        newFile.save()
        for key in strings_dict.keys():
            word = LookupWords.objects.create(
                Word=key,
                Frequency=strings_dict[key],
                File= newFile
            )
            word.save()
        print(strings_dict)
        context={
            'msg':'File uploaded Successfully.'
        }
        return render(request, 'upload.html',context)
    else:
        context={}
        return render(request, 'upload.html',context)

def profile(request):
    if request.method == "POST":
        user = User.objects.get(Username=request.session['username'])
        FirstName = request.POST.get("FirstName",'')
        LastName=request.POST.get("LastName" , "")
        Email=request.POST.get("Email" , "")
        Username=request.POST.get("Username" , "")
        Password=request.POST.get("Password" , "")
        if FirstName != '' and FirstName != ' ':
            user.FirstName = FirstName
        if LastName != '' and LastName != ' ':
            user.LastName = LastName
        if Email != '' and Email != ' ':
            user.Email = Email
        if Username != '' and Username != ' ':
            user.Username = Username
        if Password != '' and Password != ' ':
            user.Password = encrypt_password(Password)
        user.save()
        return render(request,'profile.html',{'user':user})
    else:
        user = User.objects.get(Username=request.session['username'])
        return render(request,'profile.html',{'user':user})